#!/usr/bin/env python3
"""
마크다운 파일을 PowerPoint 프레젠테이션으로 변환하는 스크립트
"""

import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor


class MarkdownToPPTX:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title, subtitle=""):
        """제목 슬라이드 추가"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(54)
        title_shape.text_frame.paragraphs[0].font.bold = True

        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(32)

        return slide

    def add_content_slide(self, title, content_items):
        """내용 슬라이드 추가"""
        if not content_items:
            return None

        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True

        # 내용
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        for item in content_items[:15]:  # 최대 15개 항목
            text = item['text']
            level = item.get('level', 0)

            p = text_frame.add_paragraph()
            p.text = text
            p.level = min(level, 2)  # 최대 레벨 2
            p.font.size = Pt(18 - level * 2)
            p.space_before = Pt(6)

            # 볼드 처리
            if item.get('bold', False):
                p.font.bold = True

        return slide

    def parse_markdown(self, md_content):
        """마크다운 파싱 및 슬라이드 생성"""
        lines = md_content.split('\n')

        main_title = None
        current_section = None
        current_subsection = None
        content_items = []

        i = 0
        while i < len(lines):
            line = lines[i].rstrip()

            # 메인 제목
            if line.startswith('# ') and not main_title:
                main_title = line[2:].strip()
                self.add_title_slide(main_title, "Network Video Recorder System")
                i += 1
                continue

            # 주요 섹션 (##)
            if line.startswith('## '):
                # 이전 서브섹션 저장
                if current_subsection and content_items:
                    slide_title = current_subsection
                    if current_section and current_section not in current_subsection:
                        slide_title = f"{current_section} - {current_subsection}"
                    self.add_content_slide(slide_title, content_items)
                    content_items = []
                # 이전 섹션 저장 (서브섹션이 없었던 경우)
                elif current_section and content_items:
                    self.add_content_slide(current_section, content_items)
                    content_items = []

                current_section = line[3:].strip()
                current_subsection = None
                i += 1
                continue

            # 서브섹션 (###)
            if line.startswith('### '):
                # 이전 서브섹션 저장
                if current_subsection and content_items:
                    slide_title = current_subsection
                    if current_section and current_section not in current_subsection:
                        slide_title = f"{current_section} - {current_subsection}"
                    self.add_content_slide(slide_title, content_items)
                    content_items = []

                current_subsection = line[4:].strip()
                i += 1
                continue

            # 서브서브섹션 (####) - 볼드 텍스트로 처리
            if line.startswith('#### '):
                text = line[5:].strip()
                content_items.append({'text': text, 'level': 0, 'bold': True})
                i += 1
                continue

            # 코드 블록 처리
            if line.startswith('```'):
                i += 1
                code_lines = []
                while i < len(lines) and not lines[i].startswith('```'):
                    code_line = lines[i].strip()
                    if code_line and not code_line.startswith('#'):
                        code_lines.append(code_line)
                    i += 1
                    if len(code_lines) >= 8:  # 최대 8줄
                        break

                # 코드 추가
                for code_line in code_lines[:8]:
                    if len(code_line) > 80:
                        code_line = code_line[:77] + "..."
                    content_items.append({'text': code_line, 'level': 1})

                i += 1
                continue

            # 리스트 항목
            if line.strip().startswith('- ') or line.strip().startswith('* '):
                # 들여쓰기 레벨 계산
                stripped = line.lstrip()
                indent = len(line) - len(stripped)
                level = indent // 2

                text = stripped[2:].strip()

                # 볼드 마커 처리 (**text**)
                bold = False
                if text.startswith('**') and '**' in text[2:]:
                    text = text.replace('**', '')
                    bold = True

                # 너무 긴 텍스트는 자르기
                if len(text) > 100:
                    text = text[:97] + "..."

                content_items.append({'text': text, 'level': level, 'bold': bold})
                i += 1
                continue

            # 숫자 리스트
            if re.match(r'^\d+\.\s', line.strip()):
                text = re.sub(r'^\d+\.\s', '', line.strip())
                # 링크 제거
                text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
                if text and len(text) < 100:
                    content_items.append({'text': text, 'level': 0})
                i += 1
                continue

            # 표 처리 (간단히)
            if '|' in line and line.strip().startswith('|'):
                # 표 구분선 건너뛰기
                if '---' in line or '━━━' in line:
                    i += 1
                    continue

                parts = [p.strip() for p in line.split('|') if p.strip()]
                if parts and len(parts) <= 3:
                    text = ' | '.join(parts)
                    if len(text) < 100:
                        content_items.append({'text': text, 'level': 0})
                i += 1
                continue

            # 일반 텍스트 (볼드 처리)
            if line.strip() and not line.startswith('---') and not line.startswith('```'):
                text = line.strip()

                # 링크 제거
                text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)

                # 볼드 마커 처리
                bold = False
                if '**' in text:
                    if text.startswith('**') and text.count('**') >= 2:
                        bold = True
                    text = text.replace('**', '')

                # 텍스트가 의미있는 경우만 추가
                if len(text) > 10 and len(text) < 150:
                    # 특수 문자로 시작하는 것은 제외
                    if not text.startswith(('┌', '│', '└', '├', '─', '·')):
                        content_items.append({'text': text, 'level': 0, 'bold': bold})

            i += 1

        # 마지막 슬라이드 저장
        if content_items:
            if current_subsection:
                slide_title = current_subsection
                if current_section and current_section not in current_subsection:
                    slide_title = f"{current_section} - {current_subsection}"
            elif current_section:
                slide_title = current_section
            else:
                slide_title = "내용"

            self.add_content_slide(slide_title, content_items)

    def convert(self, md_file, output_file):
        """마크다운 파일을 PPTX로 변환"""
        print(f"🔄 변환 중: {md_file}")

        # 마크다운 파일 읽기
        with open(md_file, 'r', encoding='utf-8') as f:
            md_content = f.read()

        # 파싱 및 슬라이드 생성
        self.parse_markdown(md_content)

        # 저장
        self.prs.save(output_file)
        print(f"✅ PowerPoint 파일 생성 완료: {output_file}")
        print(f"   총 슬라이드 수: {len(self.prs.slides)}")


def main():
    if len(sys.argv) < 2:
        print("사용법: python md_to_pptx.py <마크다운파일> [출력파일]")
        print("\n예시:")
        print("  python md_to_pptx.py user_manual.md")
        print("  python md_to_pptx.py user_manual.md output.pptx")
        sys.exit(1)

    md_file = Path(sys.argv[1])
    if not md_file.exists():
        print(f"❌ 오류: 파일을 찾을 수 없습니다 - {md_file}")
        sys.exit(1)

    # 출력 파일명 생성
    if len(sys.argv) >= 3:
        output_file = Path(sys.argv[2])
    else:
        output_file = md_file.with_suffix('.pptx')

    # 변환 실행
    converter = MarkdownToPPTX()
    converter.convert(md_file, output_file)


if __name__ == "__main__":
    main()
